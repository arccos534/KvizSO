from __future__ import annotations

import io
from datetime import datetime

from flask import Flask, jsonify, render_template, request, send_file

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Pt
except Exception:  # pragma: no cover
    Presentation = None

app = Flask(__name__)


def _to_non_negative_int(value, default: int = 0) -> int:
    try:
        return max(0, int(value))
    except (TypeError, ValueError):
        return default


def normalize_teams(raw_teams: list, *, sort_by_total: bool) -> list[dict]:
    prepared: list[dict] = []

    for i, row in enumerate(raw_teams, start=1):
        if not isinstance(row, dict):
            continue

        round1 = _to_non_negative_int(row.get("round1", 0))
        round2 = _to_non_negative_int(row.get("round2", 0))
        total = round1 + round2

        prepared.append(
            {
                "place": _to_non_negative_int(row.get("place", i), default=i),
                "team_name": str(row.get("team_name", f"Команда {i}")) or f"Команда {i}",
                "round1": round1,
                "round2": round2,
                "total": total,
            }
        )

    if sort_by_total:
        prepared.sort(key=lambda x: (-x["total"], x["place"]))
        for idx, row in enumerate(prepared, start=1):
            row["place"] = idx

    return prepared


def _set_cell_text(cell, text: str, *, bold: bool = False, align: str = "center", font_size: int = 18) -> None:
    cell.text = text
    frame = cell.text_frame
    frame.clear()
    frame.margin_left = Emu(50000)
    frame.margin_right = Emu(50000)
    frame.margin_top = Emu(25000)
    frame.margin_bottom = Emu(25000)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.word_wrap = True

    paragraph = frame.paragraphs[0]
    if align == "left":
        paragraph.alignment = PP_ALIGN.LEFT
    elif align == "right":
        paragraph.alignment = PP_ALIGN.RIGHT
    else:
        paragraph.alignment = PP_ALIGN.CENTER

    run = paragraph.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(font_size)
    run.font.name = "Calibri"
    run.font.color.rgb = RGBColor(0, 0, 0)


def build_presentation(teams: list[dict]) -> io.BytesIO:
    if Presentation is None:
        raise RuntimeError("Не установлен пакет python-pptx. Установите: pip install python-pptx")

    # Гарантируем сортировку команд по итоговым баллам перед отрисовкой в слайд.
    teams = normalize_teams(teams, sort_by_total=True)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    rows = len(teams) + 1
    cols = 5
    left = Emu(0)
    top = Emu(0)
    width = prs.slide_width
    row_height = Emu(360000)
    height = min(prs.slide_height, row_height * rows)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    table.columns[0].width = int(width * 0.09)
    table.columns[1].width = int(width * 0.33)
    table.columns[2].width = int(width * 0.17)
    table.columns[3].width = int(width * 0.17)
    table.columns[4].width = int(width * 0.19)
    for r in range(rows):
        table.rows[r].height = row_height

    headers = ["Место", "Название команды", "Раунд 1", "Раунд 2", "Результат"]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(247, 143, 111)
        _set_cell_text(cell, header, bold=True, align="left" if col == 1 else "center", font_size=20)

    for i, team in enumerate(teams, start=1):
        values = [
            str(team.get("place", i)),
            str(team.get("team_name", "")),
            str(team.get("round1", 0)),
            str(team.get("round2", 0)),
            str(team.get("total", 0)),
        ]
        for col, value in enumerate(values):
            cell = table.cell(i, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(243, 165, 142)
            _set_cell_text(cell, value, align="left" if col == 1 else "center", font_size=18)

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out


@app.get("/")
def index() -> str:
    return render_template("index.html")


@app.post("/export-pptx")
def export_pptx():
    payload = request.get_json(silent=True) or {}
    teams = payload.get("teams", [])
    if not isinstance(teams, list) or not teams:
        return jsonify({"error": "Список команд пуст."}), 400

    try:
        prepared = normalize_teams(teams, sort_by_total=True)
        if not prepared:
            return jsonify({"error": "Список команд пуст."}), 400
        pptx_stream = build_presentation(prepared)
    except RuntimeError as exc:
        return jsonify({"error": str(exc)}), 500
    except Exception:
        return jsonify({"error": "Не удалось сформировать презентацию."}), 500

    filename = f"quiz_results_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    return send_file(
        pptx_stream,
        as_attachment=True,
        download_name=filename,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


if __name__ == "__main__":
    app.run(debug=True)
